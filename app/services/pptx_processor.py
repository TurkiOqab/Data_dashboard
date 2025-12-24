"""
PowerPoint Processing Service

Extracts text, tables, and images from PPTX files.
Converts slides to images for vision analysis.
"""

import io
import os
import base64
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import json


@dataclass
class TableData:
    """Represents an extracted table."""
    rows: List[List[str]]
    headers: List[str] = field(default_factory=list)

    def to_markdown(self) -> str:
        """Convert table to markdown format."""
        if not self.rows:
            return ""

        lines = []
        headers = self.headers if self.headers else self.rows[0]
        data_rows = self.rows[1:] if not self.headers else self.rows

        # Header row
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

        # Data rows
        for row in data_rows:
            # Pad row if necessary
            padded_row = row + [""] * (len(headers) - len(row))
            lines.append("| " + " | ".join(padded_row[:len(headers)]) + " |")

        return "\n".join(lines)


@dataclass
class SlideContent:
    """Represents extracted content from a single slide."""
    slide_number: int
    title: str = ""
    text_content: List[str] = field(default_factory=list)
    tables: List[TableData] = field(default_factory=list)
    has_chart: bool = False
    has_image: bool = False
    image_path: Optional[str] = None
    image_base64: Optional[str] = None
    raw_notes: str = ""

    def get_full_text(self) -> str:
        """Get all text content as a single string."""
        parts = []
        if self.title:
            parts.append(f"Title: {self.title}")
        parts.extend(self.text_content)
        for i, table in enumerate(self.tables):
            parts.append(f"Table {i+1}:\n{table.to_markdown()}")
        if self.raw_notes:
            parts.append(f"Notes: {self.raw_notes}")
        return "\n\n".join(parts)

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for serialization."""
        return {
            "slide_number": self.slide_number,
            "title": self.title,
            "text_content": self.text_content,
            "tables": [{"rows": t.rows, "headers": t.headers} for t in self.tables],
            "has_chart": self.has_chart,
            "has_image": self.has_image,
            "image_path": self.image_path,
            "raw_notes": self.raw_notes,
            "full_text": self.get_full_text()
        }


@dataclass
class ProcessedPresentation:
    """Represents a fully processed PowerPoint file."""
    filename: str
    total_slides: int
    slides: List[SlideContent]

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for serialization."""
        return {
            "filename": self.filename,
            "total_slides": self.total_slides,
            "slides": [s.to_dict() for s in self.slides]
        }

    def save(self, output_path: str):
        """Save processed data to JSON file."""
        with open(output_path, 'w') as f:
            json.dump(self.to_dict(), f, indent=2)

    @classmethod
    def load(cls, input_path: str) -> 'ProcessedPresentation':
        """Load processed data from JSON file."""
        with open(input_path, 'r') as f:
            data = json.load(f)

        slides = []
        for s in data['slides']:
            tables = [TableData(rows=t['rows'], headers=t.get('headers', []))
                     for t in s.get('tables', [])]
            slides.append(SlideContent(
                slide_number=s['slide_number'],
                title=s.get('title', ''),
                text_content=s.get('text_content', []),
                tables=tables,
                has_chart=s.get('has_chart', False),
                has_image=s.get('has_image', False),
                image_path=s.get('image_path'),
                raw_notes=s.get('raw_notes', '')
            ))

        return cls(
            filename=data['filename'],
            total_slides=data['total_slides'],
            slides=slides
        )


class PPTXProcessor:
    """Processes PowerPoint files to extract content."""

    def __init__(self, output_dir: str = "data/processed"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def process_file(self, file_path: str) -> ProcessedPresentation:
        """Process a PowerPoint file and extract all content."""
        prs = Presentation(file_path)
        filename = Path(file_path).name
        slides = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_content = self._extract_slide_content(slide, idx, filename)
            slides.append(slide_content)

        return ProcessedPresentation(
            filename=filename,
            total_slides=len(slides),
            slides=slides
        )

    def process_uploaded_file(self, uploaded_file) -> ProcessedPresentation:
        """Process an uploaded file (Streamlit UploadedFile object)."""
        # Read the file content
        file_bytes = uploaded_file.read()
        uploaded_file.seek(0)  # Reset for potential re-reads

        # Create presentation from bytes
        prs = Presentation(io.BytesIO(file_bytes))
        filename = uploaded_file.name
        slides = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_content = self._extract_slide_content(slide, idx, filename)
            slides.append(slide_content)

        return ProcessedPresentation(
            filename=filename,
            total_slides=len(slides),
            slides=slides
        )

    def _extract_slide_content(self, slide, slide_number: int, filename: str) -> SlideContent:
        """Extract all content from a single slide."""
        content = SlideContent(slide_number=slide_number)

        # Extract title
        if slide.shapes.title:
            content.title = slide.shapes.title.text.strip()

        # Process all shapes
        for shape in slide.shapes:
            # Text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text and text != content.title:
                        content.text_content.append(text)

            # Tables
            if shape.has_table:
                table_data = self._extract_table(shape.table)
                content.tables.append(table_data)

            # Charts
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                content.has_chart = True

            # Images/Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                content.has_image = True

        # Extract notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                content.raw_notes = notes_frame.text.strip()

        return content

    def _extract_table(self, table) -> TableData:
        """Extract data from a table shape."""
        rows = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            rows.append(row_data)

        # Use first row as headers if it looks like a header
        headers = []
        if rows:
            headers = rows[0]

        return TableData(rows=rows, headers=headers)

    def slide_to_image_base64(self, file_path: str, slide_number: int) -> Optional[str]:
        """
        Convert a slide to base64 encoded image.
        Note: This requires additional dependencies or external tools.
        For now, returns None - vision analysis will work on extracted content.
        """
        # Full slide-to-image conversion requires either:
        # 1. LibreOffice/unoconv for headless conversion
        # 2. python-pptx-export (limited support)
        # 3. Aspose.Slides (commercial)
        # For the MVP, we rely on text extraction + chart detection
        return None


def extract_presentation(file_path: str) -> ProcessedPresentation:
    """Convenience function to process a presentation."""
    processor = PPTXProcessor()
    return processor.process_file(file_path)
